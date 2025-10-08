"""Access to Financial Agents Deep Dive Slide"""
import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns
import numpy as np
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

sns.set_style("whitegrid")
plt.rcParams['font.family'] = 'Arial'
output_dir = Path('efina_analysis_results')

# Load data
df = pd.read_csv(output_dir / 'prepared_dataset.csv')
top_vars = pd.read_csv(output_dir / 'top_20_variables.csv')

# ============================================================================
# VIZ 1: Components Breakdown
# ============================================================================
fig1, axes = plt.subplots(1, 3, figsize=(15, 5))

components = ['financial_agents_binary', 'transactional_account_binary', 'mobile_money_binary']
labels = ['Financial\nAgents Usage', 'Transactional\nAccount', 'Mobile\nMoney']
colors = ['#667eea', '#764ba2', '#f093fb']

for idx, (comp, label, color) in enumerate(zip(components, labels, colors)):
    # By year
    yearly_avg = df.groupby('Year')[comp].mean() * 100
    
    ax = axes[idx]
    years = yearly_avg.index
    values = yearly_avg.values
    
    bars = ax.bar(years, values, color=color, edgecolor='black', linewidth=2, width=1.5)
    
    for year, val, bar in zip(years, values, bars):
        ax.text(year, val + 2, f'{val:.1f}%', ha='center', fontsize=12, fontweight='bold')
    
    ax.set_title(label, fontsize=14, fontweight='bold')
    ax.set_ylabel('Adoption Rate (%)', fontsize=11)
    ax.set_ylim(0, 100)
    ax.set_xticks(years)
    ax.grid(True, alpha=0.3, axis='y')

plt.suptitle('Access to Financial Agents: Component Evolution (2018-2023)', 
             fontsize=16, fontweight='bold')
plt.tight_layout()
fig1_path = output_dir / 'figures' / 'agents_components.png'
plt.savefig(fig1_path, dpi=300, bbox_inches='tight', facecolor='white')
plt.close()

# ============================================================================
# VIZ 2: Access Index vs Formal Inclusion Scatter
# ============================================================================
fig2, ax = plt.subplots(figsize=(10, 7))

# Sample data for visualization (use subset for clarity)
sample = df.sample(n=min(5000, len(df)), random_state=42)

scatter = ax.scatter(sample['access_agents'], sample['FormallyIncluded'], 
                    alpha=0.4, s=20, c=sample['Year'], cmap='viridis', edgecolors='none')

# Add regression line
from scipy.stats import linregress
x = sample['access_agents'].values
y = sample['FormallyIncluded'].values
slope, intercept, r_value, p_value, std_err = linregress(x, y)
line_x = np.array([0, 1])
line_y = slope * line_x + intercept
ax.plot(line_x, line_y, 'r--', linewidth=3, label=f'R² = {r_value**2:.3f}')

ax.set_xlabel('Access to Financial Agents Index', fontsize=13, fontweight='bold')
ax.set_ylabel('Formally Included (0=No, 1=Yes)', fontsize=13, fontweight='bold')
ax.set_title('Relationship: Access to Agents vs Formal Inclusion\nStrong Positive Correlation',
             fontsize=15, fontweight='bold', pad=15)
ax.legend(fontsize=11)
ax.grid(True, alpha=0.3)

cbar = plt.colorbar(scatter, ax=ax)
cbar.set_label('Year', fontsize=11, fontweight='bold')

plt.tight_layout()
fig2_path = output_dir / 'figures' / 'agents_vs_inclusion_scatter.png'
plt.savefig(fig2_path, dpi=300, bbox_inches='tight', facecolor='white')
plt.close()

# ============================================================================
# VIZ 3: Impact Coefficient Comparison
# ============================================================================
fig3, ax = plt.subplots(figsize=(10, 8))

top_10 = top_vars.head(10).copy()
top_10 = top_10.sort_values('coefficient', ascending=True)

colors_bar = ['#ff4444' if x < 0 else '#667eea' if i == len(top_10)-1 else '#aaa' 
              for i, x in enumerate(top_10['coefficient'])]

y_pos = np.arange(len(top_10))
bars = ax.barh(y_pos, top_10['coefficient'], color=colors_bar, edgecolor='black', linewidth=1.5)

# Highlight the top one
bars[-1].set_linewidth(3)
bars[-1].set_edgecolor('#ff0000')

for i, (val, var) in enumerate(zip(top_10['coefficient'], top_10['variable'])):
    label = f'{val:.2f}'
    if val > 0:
        ax.text(val + 0.5, i, label, va='center', fontsize=11, fontweight='bold')
    else:
        ax.text(val - 0.5, i, label, va='center', ha='right', fontsize=11, fontweight='bold')

ax.set_yticks(y_pos)
ax.set_yticklabels(top_10['variable'], fontsize=11)
ax.set_xlabel('Standardized Coefficient (Impact on Inclusion)', fontsize=13, fontweight='bold')
ax.set_title('Why Access to Agents is #1: Coefficient Comparison\n(Red outline = Access to Agents with 19.78 coefficient)',
             fontsize=15, fontweight='bold', pad=15)
ax.axvline(0, color='black', linewidth=2)
ax.grid(True, alpha=0.3, axis='x')

plt.tight_layout()
fig3_path = output_dir / 'figures' / 'agents_coefficient_comparison.png'
plt.savefig(fig3_path, dpi=300, bbox_inches='tight', facecolor='white')
plt.close()

# ============================================================================
# VIZ 4: Growth Analysis
# ============================================================================
fig4, (ax1, ax2) = plt.subplots(1, 2, figsize=(14, 6))

# Left: Access index growth
yearly_access = df.groupby('Year')['access_agents'].mean()
yearly_inclusion = df.groupby('Year')['FormallyIncluded'].mean() * 100

years = yearly_access.index
access_vals = yearly_access.values

ax1.plot(years, access_vals, marker='o', linewidth=3, markersize=15, 
         color='#667eea', markerfacecolor='#764ba2', markeredgewidth=2, markeredgecolor='white')
ax1.fill_between(years, 0, access_vals, alpha=0.3, color='#667eea')

for year, val in zip(years, access_vals):
    ax1.annotate(f'{val:.3f}', (year, val), textcoords="offset points", 
                xytext=(0,10), ha='center', fontsize=13, fontweight='bold')

ax1.set_xlabel('Year', fontsize=13, fontweight='bold')
ax1.set_ylabel('Access to Agents Index', fontsize=13, fontweight='bold')
ax1.set_title('Access Index Growth\n+151% increase (2018-2023)', fontsize=14, fontweight='bold')
ax1.grid(True, alpha=0.3)
ax1.set_xticks(years)

# Right: Dual axis - Access vs Inclusion
ax2_twin = ax2.twinx()

line1 = ax2.plot(years, access_vals, marker='o', linewidth=3, markersize=12,
                 color='#667eea', label='Access Index', markeredgecolor='white', markeredgewidth=2)
line2 = ax2_twin.plot(years, yearly_inclusion.values, marker='s', linewidth=3, markersize=12,
                      color='#4caf50', label='Formal Inclusion %', markeredgecolor='white', markeredgewidth=2)

ax2.set_xlabel('Year', fontsize=13, fontweight='bold')
ax2.set_ylabel('Access to Agents Index', fontsize=12, fontweight='bold', color='#667eea')
ax2_twin.set_ylabel('Formal Inclusion Rate (%)', fontsize=12, fontweight='bold', color='#4caf50')
ax2.set_title('Parallel Growth:\nAccess & Inclusion Together', fontsize=14, fontweight='bold')
ax2.grid(True, alpha=0.3)
ax2.set_xticks(years)

# Combined legend
lines = line1 + line2
labels = [l.get_label() for l in lines]
ax2.legend(lines, labels, loc='upper left', fontsize=11)

plt.tight_layout()
fig4_path = output_dir / 'figures' / 'agents_growth_analysis.png'
plt.savefig(fig4_path, dpi=300, bbox_inches='tight', facecolor='white')
plt.close()

print("✓ All 4 visualizations created")

# ============================================================================
# Calculate key metrics
# ============================================================================

# Access change
access_2018 = df[df['Year']==2018]['access_agents'].mean()
access_2023 = df[df['Year']==2023]['access_agents'].mean()
access_change = access_2023 - access_2018
access_pct_change = (access_change / access_2018) * 100

# Inclusion change
incl_2018 = df[df['Year']==2018]['FormallyIncluded'].mean() * 100
incl_2023 = df[df['Year']==2023]['FormallyIncluded'].mean() * 100
incl_change = incl_2023 - incl_2018

# Ratio
ratio = incl_change / access_change if access_change != 0 else 0

# Coefficient
coef = top_vars[top_vars['variable']=='access_agents']['coefficient'].values[0]

# Component changes
comp_changes = {}
for comp in components:
    val_2018 = df[df['Year']==2018][comp].mean() * 100
    val_2023 = df[df['Year']==2023][comp].mean() * 100
    comp_changes[comp] = {'2018': val_2018, '2023': val_2023, 'change': val_2023 - val_2018}

print(f"\nKey Metrics:")
print(f"Access Index: {access_2018:.3f} (2018) → {access_2023:.3f} (2023)")
print(f"Change: +{access_change:.3f} ({access_pct_change:.1f}%)")
print(f"\nInclusion: {incl_2018:.1f}% (2018) → {incl_2023:.1f}% (2023)")
print(f"Change: +{incl_change:.1f}pp")
print(f"\nRatio: {ratio:.1f}pp inclusion per 1-unit access increase")
print(f"Coefficient: {coef:.2f}")

# ============================================================================
# CREATE POWERPOINT SLIDE
# ============================================================================

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])

# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.33), Inches(0.7))
title_frame = title_box.text_frame
title_frame.text = "Access to Financial Agents: The #1 Driver of Formal Inclusion"
p = title_frame.paragraphs[0]
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = RGBColor(102, 126, 234)
p.alignment = PP_ALIGN.CENTER

# Subtitle
sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(12.33), Inches(0.35))
sub_frame = sub_box.text_frame
sub_frame.text = f"Coefficient: {coef:.2f} | Growth: +{access_pct_change:.1f}% (2018-2023) | Impact: {ratio:.0f}pp per unit increase"
p = sub_frame.paragraphs[0]
p.font.size = Pt(14)
p.font.italic = True
p.font.color.rgb = RGBColor(100, 100, 100)
p.alignment = PP_ALIGN.CENTER

# Add charts
slide.shapes.add_picture(str(fig3_path), Inches(0.5), Inches(1.3), width=Inches(6))
slide.shapes.add_picture(str(fig1_path), Inches(6.8), Inches(1.3), width=Inches(6))
slide.shapes.add_picture(str(fig4_path), Inches(0.5), Inches(4.2), width=Inches(6))
slide.shapes.add_picture(str(fig2_path), Inches(6.8), Inches(4.2), width=Inches(5.8))

# Save
output_path = output_dir / 'Access_Agents_Slide.pptx'
prs.save(str(output_path))

print(f"\n✓ Slide saved: {output_path}")
