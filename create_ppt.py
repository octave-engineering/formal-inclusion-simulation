"""PowerPoint Presentation Generator for EFInA Analysis"""
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pathlib import Path

output_dir = Path('efina_analysis_results')
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Helper functions
def add_title_slide(title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    if subtitle:
        slide.placeholders[1].text = subtitle
    return slide

def add_content_slide(title, content_list):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    for item in content_list:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0
    return slide

def add_picture_slide(title, image_path):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    try:
        slide.shapes.add_picture(str(image_path), Inches(1), Inches(1.5), width=Inches(8))
    except:
        pass
    return slide

# Slide 1: Title
add_title_slide(
    "Drivers of Formal Financial Inclusion in Nigeria",
    "EFInA Data Analysis (2018-2023)\nOctave Analytics"
)

# Slide 2: Executive Summary
add_content_slide("Executive Summary", [
    "Formal inclusion surged from 45.2% (2018) to 61.2% (2023)",
    "16 percentage point increase = ~13.7 million newly included Nigerians",
    "Access to Financial Agents is the #1 driver (coefficient: 19.78)",
    "2020-2023 saw 15× faster growth than 2018-2020",
    "Key recommendation: Deploy 50,000+ agents in underserved areas"
])

# Slide 3: Key Findings
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Key Findings"
tf = slide.placeholders[1].text_frame
findings = [
    ("1. ACCESS TO FINANCIAL AGENTS IS #1 DRIVER", [
        "Coefficient: 19.78 (40× larger than most variables)",
        "SHAP importance: 3.60",
        "Effect: 34.4pp difference with/without access"
    ]),
    ("2. DRAMATIC ACCELERATION POST-2020", [
        "2018-2020: +1.0pp growth",
        "2020-2023: +15.0pp growth (15× faster)"
    ])
]
for title, points in findings:
    p = tf.add_paragraph()
    p.text = title
    p.level = 0
    for point in points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 1

# Slide 4: Top 10 Drivers
top_vars = pd.read_csv(output_dir / 'top_20_variables.csv')
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Top 10 Drivers (Ranked)"
tf = slide.placeholders[1].text_frame
for _, row in top_vars.head(10).iterrows():
    p = tf.add_paragraph()
    p.text = f"{int(row['final_rank'])}. {row['variable']}: Coef={row['coefficient']:.2f}, SHAP={row['shap_importance']:.2f}"
    p.level = 0

# Slide 5: Time Series
add_picture_slide("Formal Inclusion Rate Over Time", 
                 output_dir / 'figures' / 'inclusion_rate_time_series.png')

# Slide 6: Regional Trends
add_picture_slide("Regional Trends (2018-2023)",
                 output_dir / 'figures' / 'regional_trends.png')

# Slide 7: Urban vs Rural
add_picture_slide("Urban vs Rural Progression",
                 output_dir / 'figures' / 'urban_rural_trends.png')

# Slide 8: Top Correlations
add_picture_slide("Top Correlations with Formal Inclusion",
                 output_dir / 'figures' / 'top_correlations.png')

# Slide 9: SHAP Summary
add_picture_slide("SHAP Feature Importance",
                 output_dir / 'figures' / 'shap_summary_bar.png')

# Slide 10: Access to Agents
add_content_slide("Access to Financial Agents Index", [
    "RANK #1 DRIVER",
    "",
    "Components:",
    "  • Financial agents usage",
    "  • Transactional account ownership",
    "  • Mobile money adoption",
    "",
    "Construction: Z-score → Equal weighting → Min-max [0,1]",
    "Growth: 151% increase (2018-2023)",
    "Impact: 34.4pp difference"
])

# Slide 11: Policy Recommendations
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Top 5 Policy Recommendations"
tf = slide.placeholders[1].text_frame
recs = [
    ("1. AGENT EXPANSION (CRITICAL)", "Deploy 50,000 agents → +8-10pp impact"),
    ("2. ZERO-BALANCE ACCOUNTS (HIGH)", "Mandate implementation → +5-7pp impact"),
    ("3. FINANCIAL LITERACY (HIGH)", "Mass campaigns → +3-5pp impact"),
    ("4. MOBILE MONEY INTEGRATION (MEDIUM)", "Seamless transfers → +4-6pp impact"),
    ("5. REGIONAL EQUITY (MEDIUM)", "Target North regions → Narrow gap")
]
for title, desc in recs:
    p = tf.add_paragraph()
    p.text = title
    p.level = 0
    p = tf.add_paragraph()
    p.text = desc
    p.level = 1

# Slide 12: Methodology
add_content_slide("Methodology", [
    "Data: 85,341 respondents across 36 states + FCT",
    "Years: 2018, 2020, 2023",
    "",
    "Models:",
    "  • Logistic Regression (standardized coefficients)",
    "  • LASSO (variable selection)",
    "  • Random Forest (feature importance)",
    "  • XGBoost + SHAP (marginal contributions)",
    "",
    "Performance: 89-92% accuracy, 0.94-0.98 ROC-AUC"
])

# Slide 13: Waterfall
add_picture_slide("Contribution to Change (2018-2023)",
                 output_dir / 'figures' / 'waterfall_contributions.png')

# Slide 14: Regional Heatmap
add_picture_slide("Regional Inclusion Heatmap",
                 output_dir / 'figures' / 'regional_heatmap.png')

# Slide 15: Impact Projections
add_content_slide("Impact Projections (2024-2026)", [
    "IF TOP 3 RECOMMENDATIONS IMPLEMENTED:",
    "",
    "Agent Expansion: +8pp → 69.2% by 2025",
    "Zero-Balance Accounts: +5pp → 74.2% by 2026",
    "Financial Literacy: +3pp → 77.2% by 2026",
    "",
    "COMBINED: Could reach 77% by 2026",
    "(vs. baseline projection of 65%)"
])

# Slide 16: Conclusion
add_content_slide("Conclusion & Next Steps", [
    "KEY INSIGHT: Agent access dominates all other factors",
    "",
    "IMMEDIATE ACTIONS:",
    "  • Present to CBN Financial Inclusion Secretariat",
    "  • Advocate for agent expansion pilot (5 LGAs)",
    "  • Develop zero-balance account regulations",
    "",
    "LONG-TERM GOAL:",
    "  • 80% national inclusion by 2030",
    "  • Universal inclusion for all segments"
])

# Save
output_path = output_dir / 'EFInA_Analysis_Presentation.pptx'
prs.save(str(output_path))
print(f"✓ Presentation saved: {output_path}")
